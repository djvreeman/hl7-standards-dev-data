import argparse
import os
import sys
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import tempfile

def main():
    parser = argparse.ArgumentParser(
        description='Convert a PDF slide deck into a PowerPoint presentation with each page as an image.'
    )
    parser.add_argument('-i', '--input', required=True, help='Path to the input PDF file')
    parser.add_argument('-o', '--output', required=True, help='Path to the output PPTX file')
    parser.add_argument('-d', '--dpi', type=int, default=200, help='Resolution in DPI for the images (default: 200)')
    args = parser.parse_args()

    input_pdf = args.input
    output_pptx = args.output
    dpi = args.dpi

    # Check if input PDF exists
    if not os.path.isfile(input_pdf):
        print(f"Error: Input file '{input_pdf}' does not exist.")
        sys.exit(1)

    # Convert PDF to images
    try:
        print("Converting PDF pages to images...")
        images = convert_from_path(input_pdf, dpi=dpi)
    except Exception as e:
        print(f"Error converting PDF to images: {e}")
        sys.exit(1)

    # Create PowerPoint presentation
    prs = Presentation()

    # Set slide size to 16:9 aspect ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_slide_layout = prs.slide_layouts[6]  # Blank slide layout

    # Save the first image as the cover slide in JPEG format
    cover_slide_path = output_pptx.replace(".pptx", "-cover slide.jpg")
    try:
        print(f"Saving cover slide as '{cover_slide_path}'")
        images[0].save(cover_slide_path, 'JPEG')
    except Exception as e:
        print(f"Error saving cover slide: {e}")
        sys.exit(1)

    # Process all images and add to PowerPoint
    with tempfile.TemporaryDirectory() as tmpdirname:
        for i, image in enumerate(images):
            # Resize or crop the image to fit 16:9 aspect ratio
            image = adjust_image_to_16_9(image, dpi)

            # Save the adjusted image
            img_path = os.path.join(tmpdirname, f'slide_{i+1}.png')
            image.save(img_path, 'PNG')

            # Add image to slide
            slide = prs.slides.add_slide(blank_slide_layout)
            slide.shapes.add_picture(
                img_path, Inches(0), Inches(0),
                width=prs.slide_width, height=prs.slide_height
            )

    # Save the PowerPoint file
    try:
        prs.save(output_pptx)
        print(f"PowerPoint presentation saved to '{output_pptx}'")
    except Exception as e:
        print(f"Error saving PowerPoint file: {e}")
        sys.exit(1)

def adjust_image_to_16_9(image, dpi):
    # Calculate aspect ratios
    img_width_px, img_height_px = image.size
    img_aspect = img_width_px / img_height_px
    slide_aspect = 16 / 9

    if img_aspect > slide_aspect:
        # Image is wider than 16:9, crop sides
        new_width = int(img_height_px * slide_aspect)
        left = (img_width_px - new_width) // 2
        right = left + new_width
        top = 0
        bottom = img_height_px
    elif img_aspect < slide_aspect:
        # Image is taller than 16:9, crop top and bottom
        new_height = int(img_width_px / slide_aspect)
        top = (img_height_px - new_height) // 2
        bottom = top + new_height
        left = 0
        right = img_width_px
    else:
        # Image is already 16:9
        left, top, right, bottom = 0, 0, img_width_px, img_height_px

    # Crop the image
    image = image.crop((left, top, right, bottom))
    return image

if __name__ == '__main__':
    main()