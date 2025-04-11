#!/usr/bin/env python3
"""
Images to PowerPoint Converter

This script converts a folder of image files into a PowerPoint presentation with each image
as a separate slide. Images are automatically resized to fit the 16:9 slide format while
preserving their original aspect ratio (no cropping or stretching).

Features:
- Creates a PowerPoint presentation from a folder of images
- Supports JPG, PNG, BMP, GIF, TIFF image formats
- Supports HEIC/HEIF formats when pillow-heif package is installed
- Preserves original image aspect ratio (no cropping or stretching)
- Centers images on slides with proper letterboxing/pillarboxing
- Option to save first image as a separate cover slide file
- Option to randomize the order of images in the presentation
- Maintains image quality with configurable DPI

Requirements:
- Python 3.6+
- python-pptx
- Pillow (PIL)
- pillow-heif (optional, for HEIC support)

Installation:
    pip install python-pptx Pillow
    
    # For HEIC support:
    pip install pillow-heif

Usage:
    python images-to-pptx.py -i [input_folder] -o [output_pptx] [-d dpi] [-cover] [-randomize]

Examples:
    # Basic usage:
    python images-to-pptx.py -i ./my_images -o presentation.pptx
    
    # With specific DPI setting (higher quality):
    python images-to-pptx.py -i ./vacation_photos -o vacation.pptx -d 300
    
    # Without generating a cover slide:
    python images-to-pptx.py -i ./my_images -o presentation.pptx -cover False
    
    # Randomize the order of images:
    python images-to-pptx.py -i ./my_images -o presentation.pptx -randomize True
    
    # Combine multiple options:
    python images-to-pptx.py -i ./my_images -o presentation.pptx -d 300 -cover False -randomize True

The script will:
1. Scan the input folder for supported image files
2. Sort the images by filename (or randomize if that option is enabled)
3. Process each image to fit 16:9 slide format
4. Create a PowerPoint presentation with one slide per image
5. Save the first image as a separate file if cover slide option is enabled (default: True)
6. Output a PPTX file to the specified location

Author: Custom script
License: MIT
"""

import argparse
import os
import sys
import random
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import tempfile

# Try to import HEIC support
try:
    import pillow_heif
    pillow_heif.register_heif_opener()
    HEIC_SUPPORT = True
except ImportError:
    HEIC_SUPPORT = False

def main():
    parser = argparse.ArgumentParser(
        description='Convert a folder of images into a PowerPoint presentation with each image as a slide. Supports JPG, PNG, BMP, GIF, TIFF, and HEIC (with pillow-heif installed).'
    )
    parser.add_argument('-i', '--input', required=True, help='Path to the input folder containing images')
    parser.add_argument('-o', '--output', required=True, help='Path to the output PPTX file')
    parser.add_argument('-d', '--dpi', type=int, default=200, help='Resolution in DPI for the images (default: 200)')
    parser.add_argument('-cover', '--cover', type=lambda x: (str(x).lower() == 'true'), default=True, 
                       help='Generate a separate cover slide from the first image (default: True)')
    parser.add_argument('-randomize', '--randomize', type=lambda x: (str(x).lower() == 'true'), default=False,
                       help='Randomize the order of images in the presentation (default: False)')
    args = parser.parse_args()

    input_folder = args.input
    output_pptx = args.output
    dpi = args.dpi
    generate_cover = args.cover
    randomize_order = args.randomize

    # Check if input folder exists
    if not os.path.isdir(input_folder):
        print(f"Error: Input folder '{input_folder}' does not exist.")
        sys.exit(1)

    # Get all image files from the folder
    image_extensions = ['.jpg', '.jpeg', '.png', '.bmp', '.gif', '.tiff']
    
    # Add HEIC support if available
    if HEIC_SUPPORT:
        image_extensions.extend(['.heic', '.heif'])
    elif any(file.lower().endswith(('.heic', '.heif')) for file in os.listdir(input_folder)):
        print("WARNING: HEIC/HEIF files found but pillow-heif package is not installed.")
        print("To process HEIC files, install the package with: pip install pillow-heif")
    image_files = []
    
    for file in os.listdir(input_folder):
        file_path = os.path.join(input_folder, file)
        if os.path.isfile(file_path) and any(file.lower().endswith(ext) for ext in image_extensions):
            image_files.append(file_path)
    
    if not image_files:
        print(f"Error: No image files found in '{input_folder}'.")
        sys.exit(1)
    
    # Either sort files by name or randomize them based on the command line option
    if randomize_order:
        print(f"Randomizing the order of {len(image_files)} image files.")
        random.shuffle(image_files)
    else:
        print(f"Sorting {len(image_files)} image files by name.")
        image_files.sort()
    
    print(f"Found {len(image_files)} image files.")

    # Create PowerPoint presentation
    prs = Presentation()

    # Set slide size to 16:9 aspect ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_slide_layout = prs.slide_layouts[6]  # Blank slide layout

    # Save the first image as the cover slide in JPEG format if enabled
    if generate_cover and image_files:
        try:
            cover_slide_path = output_pptx.replace(".pptx", "-cover slide.jpg")
            print(f"Saving cover slide as '{cover_slide_path}'")
            with Image.open(image_files[0]) as img:
                # Convert to RGB if the image has an alpha channel to avoid JPEG saving issues
                if img.mode in ('RGBA', 'LA') or (img.mode == 'P' and 'transparency' in img.info):
                    rgb_img = Image.new('RGB', img.size, (255, 255, 255))
                    # Only use the image as the source if it has 3 or 4 bands
                    if len(img.getbands()) in (3, 4):
                        rgb_img.paste(img, mask=img.split()[3] if img.mode == 'RGBA' else None)
                    else:
                        rgb_img.paste(img)
                    rgb_img.save(cover_slide_path, 'JPEG')
                else:
                    img.save(cover_slide_path, 'JPEG')
        except Exception as e:
            print(f"Error saving cover slide: {e}")
            # Don't exit, just continue without saving the cover slide
            print("Continuing without saving cover slide...")
    else:
        print("Cover slide generation disabled.")


    # Process all images and add to PowerPoint
    with tempfile.TemporaryDirectory() as tmpdirname:
        for i, img_path in enumerate(image_files):
            try:
                # Open the image
                with Image.open(img_path) as image:
                    # Resize the image to fit within 16:9 aspect ratio without cropping
                    adjusted_image, position = adjust_image_to_16_9(image, dpi)
                    
                    # Save the adjusted image to a temporary file
                    temp_img_path = os.path.join(tmpdirname, f'slide_{i+1}.png')
                    adjusted_image.save(temp_img_path, 'PNG')
                    
                    # Add image to slide
                    slide = prs.slides.add_slide(blank_slide_layout)
                    
                    # Extract position values
                    left_inches, top_inches, width_inches, height_inches = position
                    
                    # Add the image with proper positioning
                    slide.shapes.add_picture(
                        temp_img_path, 
                        Inches(left_inches), 
                        Inches(top_inches),
                        width=Inches(width_inches), 
                        height=Inches(height_inches)
                    )
                    
                    print(f"Added slide {i+1} from '{os.path.basename(img_path)}'")
            except Exception as e:
                print(f"Error processing image '{img_path}': {e}")
                continue

    # Save the PowerPoint file
    try:
        prs.save(output_pptx)
        print(f"PowerPoint presentation saved to '{output_pptx}'")
    except Exception as e:
        print(f"Error saving PowerPoint file: {e}")
        sys.exit(1)

def adjust_image_to_16_9(image, dpi):
    """
    Resizes the image to fit within a 16:9 slide without cropping or distorting.
    The image will be centered on the slide with either letterboxing or pillarboxing as needed.
    
    Returns:
    - The resized image
    - The position and size to place it on the slide (left, top, width, height in inches)
    """
    # Calculate aspect ratios
    img_width_px, img_height_px = image.size
    img_aspect = img_width_px / img_height_px
    slide_aspect = 16 / 9
    
    # Standard PowerPoint slide size in inches (16:9)
    slide_width_inches = 13.333
    slide_height_inches = 7.5
    
    # Calculate scaling factors
    if img_aspect > slide_aspect:
        # Image is wider than slide, fit to width
        scale_factor = slide_width_inches / (img_width_px / dpi)
    else:
        # Image is taller than slide, fit to height
        scale_factor = slide_height_inches / (img_height_px / dpi)
    
    # Calculate new dimensions in pixels
    new_width_px = int(img_width_px * scale_factor)
    new_height_px = int(img_height_px * scale_factor)
    
    # Resize the image
    resized_image = image.resize((new_width_px, new_height_px), Image.LANCZOS)
    
    # Calculate positioning in inches to center the image on the slide
    width_inches = new_width_px / dpi
    height_inches = new_height_px / dpi
    left_inches = (slide_width_inches - width_inches) / 2
    top_inches = (slide_height_inches - height_inches) / 2
    
    return resized_image, (left_inches, top_inches, width_inches, height_inches)

if __name__ == '__main__':
    main()